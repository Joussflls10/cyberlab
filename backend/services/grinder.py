"""CyberLab Content Grinder - Full Ingestion Pipeline."""

import hashlib
import json
import asyncio
import logging
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import uuid

import fitz  # pymupdf
from pptx import Presentation
from sqlmodel import Session, select, delete

from config import get_settings
from database import engine
from .ai_client import get_client, parse_json_response

# Import models - NEVER define new tables here
from models.course import Course, Topic
from models.challenge import Challenge
from models.import_job import ImportJob

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration per spec
MAX_CHUNK_TOKENS = 6000
CHALLENGES_CACHE_DIR = Path("challenges")
MAX_CONCURRENT_JOBS = 3
MIN_APPROVED_CHALLENGES_PER_TOPIC = 3
CHALLENGE_GEN_BATCH_CONCURRENCY = 3
ENABLE_RUNTIME_ENRICHMENT = False
ENABLE_AI_REVIEW_PASSES = False


CHALLENGE_GENERATION_SYSTEM = """You are a cybersecurity instructor designing hands-on lab challenges for a Docker-based sandbox environment.

CRITICAL CONSTRAINTS — READ BEFORE GENERATING:

Each challenge runs inside a SINGLE Docker container. You cannot:
- Connect two or more containers together to simulate control/managed node topologies
- Reliably test SSH between independent hosts
- Validate workflows that require a second real machine
- Validate cloud-managed infrastructure that requires external services
- Validate workflows that require Git hosting remotes, SaaS backends, or private enterprise systems
- Simulate full virtual networking labs (hypervisor networking, multi-VM routing labs)

You CAN test:
- Writing correct file syntax/config (YAML, TOML, INI, JSON, service configs)
- Installing tools and verifying execution
- Localhost-only execution and validation
- Linting, syntax-checking, static validation
- Creating directory structures and files with required content/permissions
- Writing and executing bash/python scripts locally
- Package/service management in the local container
- Single-host command execution and output checks

CLASSIFICATION RULES:
Before generating a challenge, classify it:
- "sandboxable": can be fully tested in a single container → GENERATE IT
- "multi-host": requires two or more networked machines → SKIP IT, do not generate
- "external-dependency": requires external systems/internet/third-party platform/hypervisor integration → SKIP IT, do not generate

SKIPPING EXAMPLES:
- "SSH from Ubuntu to Rocky without password" → SKIP (multi-host)
- "Run playbook against managed node" → SKIP (multi-host)
- "Push to GitHub repo" → SKIP (external-dependency)
- "Configure VMware/VirtualBox NAT network" → SKIP (external-dependency)
- "Validate Kubernetes cluster scheduling across nodes" → SKIP (multi-host)

GENERATING EXAMPLES:
- "Write valid inventory/config files" → GENERATE
- "Write a localhost playbook and validate with ansible-playbook --syntax-check" → GENERATE
- "Write ansible.cfg or inventory and validate required keys" → GENERATE
- "Write Dockerfile and validate with dockerfile linter/syntax" → GENERATE
- "Create a systemd service file with expected directives" → GENERATE
- "Write a firewall rule file and verify syntax/tool output" → GENERATE
- "Create and validate Terraform/HCL file syntax without apply" → GENERATE

For non-sandboxable topics, DO NOT generate substitute "knowledge-only" tasks.
Instead, add explicit entries to `skipped_topics` and leave them out of `challenges`.

Never generate "write notes", "explain concept in a file", or checklist-only busywork.

Return ONLY valid JSON. No markdown, no preamble, no explanation."""


CHALLENGE_SANITY_REVIEW_SYSTEM = """You are reviewing AI-generated lab challenges for a Docker sandbox learning platform.

Your job is to REJECT bad challenges before students see them.

REJECT a challenge if:
- The command doesn't exist in a standard Linux environment
- It requires VMware, VirtualBox, or any hypervisor
- It requires a second machine or network host
- It requires GitHub or any external internet service
- It's testing a GUI action described as a CLI command
- The question is vague or just says \"run X to verify your environment\"
- The validation script just does `exit 0` or always passes
- The command is made up / not a real Linux command

APPROVE a challenge if:
- The command exists and runs in Rocky Linux / Ubuntu
- It can be fully validated inside a single container
- The question is specific and teaches something real
- The validation script actually checks meaningful state

Return JSON:
{
    \"approved\": true/false,
    \"reason\": \"why rejected or null if approved\",
    \"improved_question\": \"better version of the question or null\"
}

Return ONLY valid JSON. No markdown, no preamble, no explanation."""


QUALITY_RULES = """
QUALITY RULES — strictly enforced:
- NEVER generate "verify echo is installed" — echo is always installed, this tests nothing
- NEVER generate "create a file named topic_check.txt" — this tests nothing
- NEVER generate the same challenge twice across different topics
- Every challenge must be SPECIFIC to the topic content provided
- If the topic is "Git installation", test git specifically: git init, git config, git status
- If the topic is "Ansible setup", prioritize playbooks, ansible.cfg, inventory, syntax-check, and check-mode workflows
- Challenges must teach something real about the topic
- A student who completes your challenges should actually understand the topic better
"""


VALIDATION_PATTERNS = """
VALIDATION PATTERNS YOU CAN USE:

1. File content check:
    grep -q "expected_string" /path/to/file && exit 0 || exit 1

2. YAML syntax check:
    python3 -c "import yaml; yaml.safe_load(open('/path/file.yml'))" && exit 0 || exit 1

3. Ansible syntax check:
    ansible-playbook --syntax-check /path/playbook.yml > /dev/null 2>&1 && exit 0 || exit 1

4. Check FQCN usage in playbook:
    grep -q "ansible.builtin." /path/playbook.yml && exit 0 || exit 1

5. Check become is used:
    grep -q "become: true" /path/playbook.yml && exit 0 || exit 1

6. Check inventory structure:
    grep -q "ansible_host:" /path/hosts.yml && exit 0 || exit 1

7. Check ansible.cfg has correct key:
    grep -q "inventory" /path/ansible.cfg && exit 0 || exit 1

8. Check file is executable:
    test -x /path/script.sh && exit 0 || exit 1

9. Check directory structure exists:
    test -d /path/playbooks && test -d /path/inventories && exit 0 || exit 1

10. Run playbook in check mode:
     ansible-playbook --check /path/playbook.yml -i localhost, 2>&1 | grep -q "0 failed" && exit 0 || exit 1
"""


CHALLENGE_TYPES_REQUIRED = """
For topics involving Ansible, you MUST generate at least one challenge of each type:
- Write a playbook file with specific requirements (use YAML syntax validation)
- Write correct ansible.cfg or inventory.yml (use grep to check required keys)
- Run a real ansible command and check meaningful output (not just --version)

EXAMPLE of a GOOD challenge for "Ansible Bootstrap Playbook":
{
  "question": "Write a playbook at /root/bootstrap.yml that installs 'nmap' on localhost using ansible.builtin.dnf with FQCN and become: true",
  "type": "file",
  "sandbox_image": "rocky9-base",
  "difficulty": "medium",
  "validation_script": "#!/bin/bash\\ntest -f /root/bootstrap.yml || exit 1\\ngrep -q 'ansible.builtin.dnf' /root/bootstrap.yml || exit 1\\ngrep -q 'become: true' /root/bootstrap.yml || exit 1\\npython3 -c \"import yaml; yaml.safe_load(open('/root/bootstrap.yml'))\" || exit 1\\nansible-playbook --syntax-check /root/bootstrap.yml -i localhost, > /dev/null 2>&1 || exit 1\\nexit 0"
}

EXAMPLE of a BAD challenge (never generate this):
{
  "question": "Run ansible --version and save output to ansible_version.txt",
  "validation_script": "test -f ansible_version.txt && exit 0 || exit 1"
}
"""


ENRICHMENT_SYSTEM = """You are a cybersecurity and infrastructure expert converting a student task sheet into rich educational content.

Your job is NOT to summarize. Your job is to EXPAND and EXPLAIN.

For every task or command mentioned, add:
- What it does and why it exists
- What the correct syntax looks like with a real example
- What a common mistake looks like and why it fails
- What success looks like (expected output or state)
- How it connects to the broader topic

Transform bullet points into paragraphs. Transform task instructions into knowledge.
A student who reads your output should understand the topic deeply, not just know what to type.

Write in English regardless of input language.

Return plain educational text only. No markdown fences."""

GENERIC_LOW_VALUE_COMMANDS = {
    "echo",
    "cat",
    "ls",
    "pwd",
    "cd",
    "touch",
    "mkdir",
    "rm",
    "true",
    "false",
}

NON_SANDBOXABLE_HINT_TERMS = {
    "snapshot",
    "vmware",
    "virtualbox",
    "vcenter",
    "esxi",
    "ova",
    "ovf",
    "vmnet",
}

DOMAIN_PROFILES: Dict[str, Dict[str, Any]] = {
    "ansible": {
        "keywords": [
            "ansible",
            "playbook",
            "ansible.cfg",
            "inventory",
            "ansible-playbook",
            "ansible.builtin",
        ],
    },
    "git": {
        "keywords": [
            "git",
            "repository",
            "commit",
            "branch",
            "version control",
            "git config",
            "git init",
        ],
    },
    "linux": {
        "keywords": [
            "linux",
            "bash",
            "shell",
            "filesystem",
            "permissions",
            "process",
            "systemd",
            "package",
            "chmod",
            "chown",
        ],
    },
}

DOMAIN_CHALLENGE_REQUIREMENTS: Dict[str, str] = {
    "ansible": CHALLENGE_TYPES_REQUIRED,
    "git": """
For topics involving Git, you MUST generate at least one challenge of each type:
- Repository state challenge (init/status/repo structure)
- Commit/config challenge (local user/email and meaningful commit)
- Branch/history challenge (branch switch/log verification)

BAD EXAMPLES (never generate as primary learning tasks):
- "Run git --version and save output"
- "Run git --help and capture output"
- "Write notes about Git concepts"
""",
    "linux": """
For Linux topics, prioritize practical host tasks:
- Filesystem/structure task (directories/files with expected content)
- Permission/executable task (chmod/stat-based verification)
- Execution task (script or command output verified by content checks)

Avoid vague prompts like "verify your environment".
""",
}

DOMAIN_VALIDATION_PATTERNS: Dict[str, str] = {
    "git": """
GIT VALIDATION EXAMPLES:
- test -d /root/git-lab/.git || exit 1
- git -C /root/git-lab config --local user.email | grep -q "@" || exit 1
- git -C /root/git-lab log --oneline | grep -q "Initial" || exit 1
- git -C /root/git-lab branch --show-current | grep -q "feature/" || exit 1
""",
    "linux": """
LINUX VALIDATION EXAMPLES:
- test -d /root/linux-lab/logs || exit 1
- stat -c '%a' /root/linux-lab/report.txt | grep -q '^640$' || exit 1
- test -x /root/linux-lab/check_disk.sh || exit 1
- grep -q "Filesystem" /root/linux-lab/disk.txt || exit 1
""",
}

# Global semaphore to limit concurrent processing
_job_semaphore = asyncio.Semaphore(MAX_CONCURRENT_JOBS)


def _fallback_extract_topics(text: str) -> Dict[str, Any]:
    """Deterministic fallback topic extraction when AI is unavailable."""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    course_name = lines[0][:120] if lines else "Imported Course"

    candidates: List[str] = []

    # Prefer numbered heading-like lines (e.g., "2.8 Ansible installeren ...").
    numbered_titles = re.findall(r"(?:^|\n)\s*\d+(?:\.\d+)*\s+[\-–:]?\s*([^\n]{8,160})", text)
    for title in numbered_titles:
        clean = re.sub(r"\s+", " ", title).strip(" -:\t")
        if clean:
            candidates.append(clean)

    # Then use paragraph opening lines.
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    for paragraph in paragraphs:
        first_line = paragraph.splitlines()[0].strip()
        if first_line:
            candidates.append(first_line[:120])

    # If still sparse, use non-empty lines that look like section/task titles.
    for line in lines:
        if len(line) < 8 or len(line) > 160:
            continue
        if re.search(r"[.!?]", line) and len(line.split()) > 16:
            continue
        candidates.append(line)

    candidates = _unique_preserve_order(candidates)
    if not candidates:
        candidates = ["Core Concepts"]
    candidates = candidates[:10]

    topics = []
    for idx, name in enumerate(candidates, start=1):
        topics.append(
            {
                "name": name,
                "order": idx,
                "key_concepts": [name],
                "tools_mentioned": [],
                "commands_mentioned": [],
                "procedures": [],
            }
        )

    return {
        "course_name": course_name,
        "topics": topics,
    }


def _topic_slug(name: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", (name or "topic").strip().lower()).strip("-")
    return slug[:40] or "topic"


def _normalize_topic_key(name: str) -> str:
    parts = [p for p in re.split(r"[^a-z0-9]+", (name or "").lower()) if p]
    return "-".join(parts[:12]) or _topic_slug(name)


def _unique_preserve_order(items: List[str]) -> List[str]:
    seen: set[str] = set()
    out: List[str] = []
    for item in items:
        text = str(item or "").strip()
        if not text:
            continue
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(text)
    return out


def _merge_topic_payload(existing: Dict[str, Any], incoming: Dict[str, Any]) -> Dict[str, Any]:
    for field in ["key_concepts", "tools_mentioned", "commands_mentioned", "procedures"]:
        existing_vals = existing.get(field, []) if isinstance(existing.get(field), list) else []
        incoming_vals = incoming.get(field, []) if isinstance(incoming.get(field), list) else []
        existing[field] = _unique_preserve_order(existing_vals + incoming_vals)
    return existing


def _topic_corpus(topic_data: Dict[str, Any]) -> str:
    corpus_parts = [str(topic_data.get("name") or "")]
    for field in ["key_concepts", "tools_mentioned", "commands_mentioned", "procedures"]:
        values = topic_data.get(field, [])
        if isinstance(values, list):
            corpus_parts.extend(str(v or "") for v in values)
    return " ".join(corpus_parts).lower()


def _keyword_in_corpus(corpus: str, keyword: str) -> bool:
    needle = (keyword or "").strip().lower()
    if not needle:
        return False

    # Phrase-like keywords are best matched by substring.
    if re.search(r"[^a-z0-9]", needle):
        return needle in corpus

    # Single-token keywords should use a boundary-aware check.
    return bool(re.search(rf"\b{re.escape(needle)}\b", corpus))


def _detect_topic_domain(topic_data: Dict[str, Any]) -> str:
    corpus = _topic_corpus(topic_data)
    best_domain = "generic"
    best_score = 0

    for domain, profile in DOMAIN_PROFILES.items():
        keywords = profile.get("keywords", []) if isinstance(profile, dict) else []
        score = sum(1 for kw in keywords if _keyword_in_corpus(corpus, str(kw)))
        if score > best_score:
            best_domain = domain
            best_score = score

    return best_domain


def _build_validation_patterns(topic_data: Dict[str, Any]) -> str:
    domain = _detect_topic_domain(topic_data)
    domain_patterns = DOMAIN_VALIDATION_PATTERNS.get(domain)
    if domain_patterns:
        return f"{VALIDATION_PATTERNS}\n\n{domain_patterns.strip()}"
    return VALIDATION_PATTERNS


def _build_challenge_type_requirements(topic_data: Dict[str, Any]) -> str:
    domain = _detect_topic_domain(topic_data)
    return DOMAIN_CHALLENGE_REQUIREMENTS.get(domain, "")


def _classify_topic_sandboxability(topic_data: Dict[str, Any]) -> Optional[str]:
    """Return a reason if this topic appears non-sandboxable in a single container."""
    corpus = _topic_corpus(topic_data)

    hypervisor_terms = ["vmware", "virtualbox", "vcenter", "esxi", "vmnet", "ovf", "ova", "snapshot", "hyper-v"]
    if any(term in corpus for term in hypervisor_terms):
        return "external-dependency: hypervisor/VM workflow cannot be validated inside a single container"

    external_terms = ["github", "gitlab", "bitbucket", "pull request", "remote repository"]
    if any(term in corpus for term in external_terms):
        return "external-dependency: requires external hosted services"

    multihost_terms = ["managed node", "remote ssh", "second host", "another machine", "multi-host"]
    if any(term in corpus for term in multihost_terms):
        return "multi-host: requires more than one machine/network host"

    return None


def _extract_candidate_commands(topic_data: Dict[str, Any]) -> List[str]:
    commands = topic_data.get("commands_mentioned", []) or []
    candidates: List[str] = []
    for raw_cmd in commands:
        cmd = str(raw_cmd or "").strip()
        if not cmd:
            continue
        token = cmd.split()[0].strip().lower()
        if not re.fullmatch(r"[a-z0-9._+\-]+", token):
            continue
        if token in GENERIC_LOW_VALUE_COMMANDS:
            continue
        if token in NON_SANDBOXABLE_HINT_TERMS:
            continue
        if token in candidates:
            continue
        candidates.append(token)
    return candidates


def _challenge_fingerprint(question: str, validation_script: str) -> str:
    normalized_q = re.sub(r"\s+", " ", (question or "").strip().lower())
    normalized_s = re.sub(r"\s+", " ", (validation_script or "").strip().lower())
    return hashlib.sha256(f"{normalized_q}::{normalized_s}".encode()).hexdigest()


def _is_generic_filler_challenge(question: str, validation_script: str) -> Optional[str]:
    q = (question or "").strip().lower()
    s = (validation_script or "").strip().lower()

    if "verify `echo`" in q or "verify echo" in q:
        return "generic low-value echo verification"
    if "topic_check.txt" in q or "topic_check.txt" in s:
        return "generic topic_check template"
    if "verify your environment" in q:
        return "vague environment verification"
    if re.search(r"create\s+a\s+file\s+named\s+`?topic", q):
        return "generic topic marker file challenge"
    if "command -v echo" in s:
        return "validation checks only echo availability"
    if any(term in q for term in ["notes.md", "concepts.txt", "documenting how", "explain how the concept", "checklist_"]):
        return "knowledge-only filler task"
    return None


def _is_ansible_topic(topic_data: Dict[str, Any]) -> bool:
    return _detect_topic_domain(topic_data) == "ansible"


def _is_low_value_domain_challenge(question: str, validation_script: str, topic_data: Dict[str, Any]) -> Optional[str]:
    domain = _detect_topic_domain(topic_data)
    if domain == "generic":
        return None

    q = (question or "").strip().lower()
    s = (validation_script or "").strip().lower()

    knowledge_only_terms = ["notes.md", "concept", "document", "checklist", "explain"]
    if any(term in q for term in knowledge_only_terms):
        return f"{domain} topic: knowledge-only writing challenge is low-value"

    version_or_help_only = any(token in q for token in ["--version", " --help", "-help"])
    output_dump = "save output" in q or "capture output" in q

    if domain == "ansible":
        if "ansible" in q and version_or_help_only:
            mentions_real_work = any(token in q for token in ["playbook", "ansible.cfg", "inventory", "syntax-check", "--check"])
            if not mentions_real_work:
                return "ansible topic: version/help-only challenge is low-value"

        if "ansible" in q and output_dump and any(token in q for token in ["_version.txt", "_help.txt"]):
            return "ansible topic: output-dump challenge is low-value"

        if "ansible" in s and "--version" in s and "--syntax-check" not in s and "--check" not in s:
            return "ansible validation: version-only checks are low-value"

        return None

    if domain == "git":
        if "git" in q and version_or_help_only:
            return "git topic: version/help-only challenge is low-value"

        if "git" in q and output_dump and any(token in q for token in ["_version.txt", "_help.txt"]):
            return "git topic: output-dump challenge is low-value"

        if "command -v git" in s and all(token not in s for token in [".git", "git log", "git config", "git branch", "git commit"]):
            return "git validation: command-exists-only checks are low-value"

        return None

    if domain == "linux":
        if version_or_help_only and output_dump:
            return f"{domain} topic: help/version output capture is low-value"

        if "verify your environment" in q:
            return f"{domain} topic: vague environment verification is low-value"

        if "command -v" in s and all(token not in s for token in ["grep -q", "stat -c", "test -", "ss -", "ip "]):
            return f"{domain} validation: command-exists-only checks are low-value"

        return None

    return None


def _is_low_value_ansible_challenge(question: str, validation_script: str, topic_data: Dict[str, Any]) -> Optional[str]:
    if not _is_ansible_topic(topic_data):
        return None
    return _is_low_value_domain_challenge(question, validation_script, topic_data)


def _is_weak_validation_script(validation_script: str) -> Optional[str]:
    script = (validation_script or "").strip()
    if not script:
        return "missing validation script"

    first_line = script.splitlines()[0].strip() if script.splitlines() else ""
    if first_line.startswith("#!"):
        if not re.fullmatch(r"^#!\s*(/bin/(ba)?sh|/usr/bin/env\s+bash)\s*$", first_line):
            return "invalid or unsupported shebang"

    lower = script.lower()

    if re.fullmatch(r"(?is)\s*(#!/bin/(ba)?sh\s*)?exit\s+0\s*;?\s*", script):
        return "validation script always passes"

    lines = [ln.strip() for ln in script.splitlines() if ln.strip() and not ln.strip().startswith("#")]
    if not lines:
        return "validation script has no executable checks"

    # Detect files-only checks that never verify content/syntax/behavior.
    non_terminal = [ln for ln in lines if ln not in {"exit 0", "exit 1"}]
    if non_terminal and all(re.search(r"^\[\s*-f\s+[^\]]+\]\s*(\|\|\s*exit\s+1)?$", ln) for ln in non_terminal):
        return "validation script checks file existence only"

    # Detect pipeline checks whose status is ignored due to final unconditional exit 0.
    if lines and lines[-1] == "exit 0":
        for ln in lines[:-1]:
            if "|" in ln and "|| exit 1" not in ln and "&&" not in ln and not ln.startswith("if "):
                return "validation script may ignore pipeline failure before unconditional exit 0"

    return None


def _validate_ansible_challenge_mix(challenges: List[Dict[str, Any]], topic_data: Dict[str, Any]) -> Optional[str]:
    if not _is_ansible_topic(topic_data):
        return None

    has_playbook = False
    has_config_or_inventory = False
    has_meaningful_execution = False

    for challenge in challenges:
        q = str(challenge.get("question") or "").lower()
        s = str(challenge.get("validation_script") or "").lower()

        if any(token in q for token in ["playbook", ".yml", ".yaml"]) or "ansible-playbook" in s:
            has_playbook = True

        if any(token in q for token in ["ansible.cfg", "inventory", "hosts.yml", "hosts.ini"]):
            has_config_or_inventory = True
        if any(token in s for token in ["ansible_host:", "inventory", "[defaults]", "ansible.cfg"]):
            has_config_or_inventory = True

        low_value = any(token in q for token in ["--version", "--help", "_version.txt", "_help.txt"])
        meaningful_markers = any(token in q for token in ["syntax-check", "--check", "localhost", "ansible-config", "playbook", "inventory", "ansible.cfg"])
        meaningful_markers = meaningful_markers or any(token in s for token in ["--syntax-check", "--check", "ansible-config", "ansible.builtin."])
        if not low_value and meaningful_markers:
            has_meaningful_execution = True

    missing = []
    if not has_playbook:
        missing.append("playbook")
    if not has_config_or_inventory:
        missing.append("ansible.cfg/inventory")
    if not has_meaningful_execution:
        missing.append("meaningful ansible execution")

    if missing:
        return f"ansible challenge mix missing required types: {', '.join(missing)}"
    return None


def _validate_git_challenge_mix(challenges: List[Dict[str, Any]]) -> Optional[str]:
    has_repo_state = False
    has_config_or_commit = False
    has_branch_or_history = False

    for challenge in challenges:
        q = str(challenge.get("question") or "").lower()
        s = str(challenge.get("validation_script") or "").lower()

        if any(token in q for token in ["git init", "initialize a git repository", ".git"]) or ".git" in s:
            has_repo_state = True

        if any(token in q for token in ["git config", "commit", "user.email", "user.name"]):
            has_config_or_commit = True
        if any(token in s for token in ["git config", "git log", "git commit"]):
            has_config_or_commit = True

        if any(token in q for token in ["branch", "checkout", "switch", "log", "history"]):
            has_branch_or_history = True
        if any(token in s for token in ["git branch", "--show-current", "git log"]):
            has_branch_or_history = True

    missing = []
    if not has_repo_state:
        missing.append("repository state/init")
    if not has_config_or_commit:
        missing.append("config/commit")
    if not has_branch_or_history:
        missing.append("branch/history")

    if missing:
        return f"git challenge mix missing required types: {', '.join(missing)}"
    return None


def _validate_domain_challenge_mix(challenges: List[Dict[str, Any]], topic_data: Dict[str, Any]) -> Optional[str]:
    domain = _detect_topic_domain(topic_data)
    if domain == "ansible":
        return _validate_ansible_challenge_mix(challenges, topic_data)
    if domain == "git":
        return _validate_git_challenge_mix(challenges)
    return None


def _fallback_generate_challenges(topic_data: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Fast deterministic fallback that avoids knowledge-only slop."""
    topic_name = str(topic_data.get("name") or "Unknown Topic")
    skip_reason = _classify_topic_sandboxability(topic_data)
    if skip_reason:
        return []

    domain = _detect_topic_domain(topic_data)

    if domain == "ansible":
        return [
            {
                "question": "Write a playbook at /root/bootstrap.yml that installs 'nmap' on localhost using ansible.builtin.dnf with become: true.",
                "hint": "Use YAML with hosts: localhost and a tasks list using ansible.builtin.dnf.",
                "type": "file",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -f /root/bootstrap.yml || exit 1\n"
                    "grep -q 'ansible.builtin.dnf' /root/bootstrap.yml || exit 1\n"
                    "grep -q 'become: true' /root/bootstrap.yml || exit 1\n"
                    "python3 -c \"import yaml; yaml.safe_load(open('/root/bootstrap.yml'))\" || exit 1\n"
                    "ansible-playbook --syntax-check /root/bootstrap.yml -i localhost, > /dev/null 2>&1 || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Create /root/ansible.cfg with a [defaults] section that sets inventory=/root/inventory.yml and host_key_checking=False.",
                "hint": "Use INI format and include both required keys.",
                "type": "file",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -f /root/ansible.cfg || exit 1\n"
                    "grep -q '^\\[defaults\\]' /root/ansible.cfg || exit 1\n"
                    "grep -q '^inventory\\s*=\\s*/root/inventory.yml' /root/ansible.cfg || exit 1\n"
                    "grep -q '^host_key_checking\\s*=\\s*False' /root/ansible.cfg || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Create /root/inventory.yml with localhost and ansible_host, then run /root/bootstrap.yml in check mode and save output to /root/check_mode_output.txt.",
                "hint": "Ensure inventory contains localhost with ansible_host and use ansible-playbook --check.",
                "type": "command",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -f /root/inventory.yml || exit 1\n"
                    "grep -q 'ansible_host:' /root/inventory.yml || exit 1\n"
                    "test -f /root/bootstrap.yml || exit 1\n"
                    "ansible-playbook --check /root/bootstrap.yml -i /root/inventory.yml > /root/check_mode_output.txt 2>&1 || exit 1\n"
                    "grep -Eq 'failed=0|0 failed' /root/check_mode_output.txt || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
        ]

    if domain == "git":
        return [
            {
                "question": "Initialize a Git repository at /root/git-lab, create README.md with text 'CyberLab Git Lab', and make an initial commit with message 'Initial commit'.",
                "hint": "Use git init, add, and commit in /root/git-lab.",
                "type": "command",
                "sandbox_image": "rocky9-base",
                "difficulty": "easy",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -d /root/git-lab/.git || exit 1\n"
                    "grep -q 'CyberLab Git Lab' /root/git-lab/README.md || exit 1\n"
                    "git -C /root/git-lab log --oneline | grep -q 'Initial commit' || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Configure local Git identity in /root/git-lab with user.name='Cyber Student' and user.email='student@example.com'.",
                "hint": "Set local config values with git config --local.",
                "type": "command",
                "sandbox_image": "rocky9-base",
                "difficulty": "easy",
                "validation_script": (
                    "#!/bin/bash\n"
                    "git -C /root/git-lab config --local user.name | grep -q '^Cyber Student$' || exit 1\n"
                    "git -C /root/git-lab config --local user.email | grep -q '^student@example.com$' || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Create and switch to branch feature/hardening in /root/git-lab, add hardening.txt containing 'enabled', and commit with message 'Add hardening notes'.",
                "hint": "Use git switch -c (or checkout -b), then add and commit.",
                "type": "command",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "git -C /root/git-lab branch --show-current | grep -q '^feature/hardening$' || exit 1\n"
                    "git -C /root/git-lab show --name-only --oneline HEAD | grep -q 'hardening.txt' || exit 1\n"
                    "git -C /root/git-lab log -1 --pretty=%s | grep -q '^Add hardening notes$' || exit 1\n"
                    "grep -q '^enabled$' /root/git-lab/hardening.txt || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
        ]

    if domain == "linux":
        return [
            {
                "question": "Create /root/linux-lab with subdirectory /root/linux-lab/logs and a file /root/linux-lab/report.txt containing 'linux-lab-ready'.",
                "hint": "Use mkdir -p and output redirection.",
                "type": "file",
                "sandbox_image": "rocky9-base",
                "difficulty": "easy",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -d /root/linux-lab/logs || exit 1\n"
                    "test -f /root/linux-lab/report.txt || exit 1\n"
                    "grep -q '^linux-lab-ready$' /root/linux-lab/report.txt || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Set permissions so /root/linux-lab/report.txt is 640 and /root/linux-lab/logs is 750.",
                "hint": "Use chmod and verify with stat.",
                "type": "command",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "stat -c '%a' /root/linux-lab/report.txt | grep -q '^640$' || exit 1\n"
                    "stat -c '%a' /root/linux-lab/logs | grep -q '^750$' || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
            {
                "question": "Create an executable script /root/linux-lab/check_disk.sh that writes `df -h` output to /root/linux-lab/disk.txt and run it once.",
                "hint": "Include a shebang and `df -h > /root/linux-lab/disk.txt` in the script.",
                "type": "file",
                "sandbox_image": "rocky9-base",
                "difficulty": "medium",
                "validation_script": (
                    "#!/bin/bash\n"
                    "test -f /root/linux-lab/check_disk.sh || exit 1\n"
                    "test -x /root/linux-lab/check_disk.sh || exit 1\n"
                    "grep -q 'df -h' /root/linux-lab/check_disk.sh || exit 1\n"
                    "test -s /root/linux-lab/disk.txt || exit 1\n"
                    "grep -q 'Filesystem' /root/linux-lab/disk.txt || exit 1\n"
                    "exit 0"
                ),
                "expected_output": None,
            },
        ]

    topic_slug = _topic_slug(topic_name)
    commands = _extract_candidate_commands(topic_data)
    if not commands:
        tools = [str(t).strip().lower() for t in (topic_data.get("tools_mentioned", []) or []) if str(t).strip()]
        for tool in tools:
            token = re.sub(r"[^a-z0-9._+\-]", "", tool)
            if token and token not in GENERIC_LOW_VALUE_COMMANDS and token not in NON_SANDBOXABLE_HINT_TERMS:
                commands.append(token)
    commands = _unique_preserve_order(commands)
    if not commands:
        return []

    primary = commands[0]
    primary_var = re.sub(r"[^a-z0-9_]+", "_", primary)
    verify_script = f"/root/verify_{primary_var}.sh"
    verify_log = f"/root/{primary_var}_verify.log"
    command_path_file = f"/root/{primary_var}_path.txt"
    workspace = f"/root/{topic_slug}-workspace"

    return [
        {
            "question": f"Create {verify_script} that checks `{primary}` exists (using `command -v`) and writes a short success message before exiting 0.",
            "hint": "Add a shebang, use command -v, and make the script executable.",
            "type": "file",
            "sandbox_image": "rocky9-base",
            "difficulty": "easy",
            "validation_script": (
                "#!/bin/bash\n"
                f"test -f {verify_script} || exit 1\n"
                f"test -x {verify_script} || exit 1\n"
                f"grep -q 'command -v {primary}' {verify_script} || exit 1\n"
                "exit 0"
            ),
            "expected_output": None,
        },
        {
            "question": f"Run {verify_script} and save its output to {verify_log}.",
            "hint": "Execute the script and redirect stdout/stderr to the log file.",
            "type": "command",
            "sandbox_image": "rocky9-base",
            "difficulty": "medium",
            "validation_script": (
                "#!/bin/bash\n"
                f"test -x {verify_script} || exit 1\n"
                f"test -s {verify_log} || exit 1\n"
                "exit 0"
            ),
            "expected_output": None,
        },
        {
            "question": f"Create directory {workspace} and save the absolute path of `{primary}` into {command_path_file}.",
            "hint": "Use `mkdir -p` and `command -v` to capture the command path.",
            "type": "file",
            "sandbox_image": "rocky9-base",
            "difficulty": "medium",
            "validation_script": (
                "#!/bin/bash\n"
                f"test -d {workspace} || exit 1\n"
                f"test -s {command_path_file} || exit 1\n"
                f"grep -Eq '^/' {command_path_file} || exit 1\n"
                f"grep -q '{primary}' {command_path_file} || exit 1\n"
                "exit 0"
            ),
            "expected_output": None,
        },
    ]


def compute_source_hash(file_path: str) -> str:
    """Compute SHA256 hash of file for duplicate detection."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def compute_challenge_id(question: str, validation_script: str, topic_id: str) -> str:
    """Compute SHA256 hash for challenge deduplication per spec."""
    content = f"{question}{validation_script}{topic_id}"
    return hashlib.sha256(content.encode()).hexdigest()


def parse_pdf(file_path: str) -> str:
    """Extract text from PDF using pymupdf."""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        logger.error(f"Failed to parse PDF {file_path}: {e}")
        raise
    return text


def parse_pptx(file_path: str) -> str:
    """Extract text from PowerPoint using python-pptx."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.has_notes_slide:
                    text += shape.notes_slide.notes_text_frame.text + "\n"
    except Exception as e:
        logger.error(f"Failed to parse PPTX {file_path}: {e}")
        raise
    return text


def chunk_text(text: str, max_tokens: int = 6000) -> List[str]:
    """Chunk text at ~6000 tokens (approx 24000 chars)."""
    max_chars = max_tokens * 4
    chunks = []
    current_chunk = ""
    
    for paragraph in text.split("\n\n"):
        if len(current_chunk) + len(paragraph) > max_chars:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = paragraph
        else:
            current_chunk += "\n\n" + paragraph
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks if chunks else [text]


async def extract_topics(text: str) -> Dict[str, Any]:
    """Extract topics from course text using DeepSeek V3."""
    client = get_client()
    
    system_prompt = """You are a curriculum analyst. Extract structured learning topics from course material.
Return ONLY valid JSON. No markdown, no preamble, no explanation."""
    
    chunks = chunk_text(text, max_tokens=MAX_CHUNK_TOKENS)
    topic_map: Dict[str, Dict[str, Any]] = {}
    course_name: Optional[str] = None
    chunks_with_topics = 0
    chunk_failures = 0

    for chunk_idx, chunk in enumerate(chunks, start=1):
        user_prompt = f"""Analyze this course material chunk and extract distinct topics and subtopics.
For each topic, list key commands, tools, and procedures covered in THIS CHUNK.

Chunk context: {chunk_idx}/{len(chunks)}

Return this exact JSON structure:
{{
 "course_name": "inferred name of the course",
 "topics": [
   {{
     "name": "Topic name",
     "order": 1,
     "key_concepts": ["concept1", "concept2"],
     "tools_mentioned": ["ansible", "wazuh"],
     "commands_mentioned": ["ansible-playbook", "systemctl"],
     "procedures": ["step-by-step procedure descriptions extracted verbatim"]
   }}
 ]
}}

Course material chunk:
{chunk}"""

        try:
            response = await client.call_model("grinder", system_prompt, user_prompt, max_tokens=4000)
            parsed = parse_json_response(response)
            parsed_topics = parsed.get("topics", []) if isinstance(parsed.get("topics", []), list) else []

            maybe_course_name = str(parsed.get("course_name") or "").strip()
            if maybe_course_name and not course_name:
                course_name = maybe_course_name

            if parsed_topics:
                chunks_with_topics += 1

            for pos, topic in enumerate(parsed_topics, start=1):
                if not isinstance(topic, dict):
                    continue
                name = str(topic.get("name") or "").strip()
                if not name:
                    continue

                key = _normalize_topic_key(name)
                incoming = {
                    "name": name,
                    "order": topic.get("order", pos),
                    "key_concepts": topic.get("key_concepts", []) if isinstance(topic.get("key_concepts", []), list) else [],
                    "tools_mentioned": topic.get("tools_mentioned", []) if isinstance(topic.get("tools_mentioned", []), list) else [],
                    "commands_mentioned": topic.get("commands_mentioned", []) if isinstance(topic.get("commands_mentioned", []), list) else [],
                    "procedures": topic.get("procedures", []) if isinstance(topic.get("procedures", []), list) else [],
                    "_first_chunk": chunk_idx,
                    "_first_pos": pos,
                }

                existing = topic_map.get(key)
                if not existing:
                    topic_map[key] = incoming
                else:
                    merged = _merge_topic_payload(existing, incoming)
                    topic_map[key] = merged
        except Exception as e:
            if 'rate-limited' in str(e).lower() or '429' in str(e).lower():
                raise
            chunk_failures += 1
            logger.warning("Topic extraction failed for chunk %s/%s: %s", chunk_idx, len(chunks), e)

    if not topic_map:
        logger.warning("AI topic extraction returned empty topics; using fallback extraction")
        fallback = _fallback_extract_topics(text)
        fallback["_generation_mode"] = "fallback"
        fallback["_fallback_reason"] = (
            f"No topics extracted across {len(chunks)} chunks (chunk_failures={chunk_failures})"
        )
        fallback["_chunk_count"] = len(chunks)
        fallback["_chunks_with_topics"] = chunks_with_topics
        fallback["_chunk_failures"] = chunk_failures
        return fallback

    ordered = sorted(topic_map.values(), key=lambda t: (t.get("_first_chunk", 1), t.get("_first_pos", 1)))
    final_topics: List[Dict[str, Any]] = []
    for idx, topic in enumerate(ordered, start=1):
        final_topics.append(
            {
                "name": topic.get("name", f"Topic {idx}"),
                "order": idx,
                "key_concepts": _unique_preserve_order(topic.get("key_concepts", [])),
                "tools_mentioned": _unique_preserve_order(topic.get("tools_mentioned", [])),
                "commands_mentioned": _unique_preserve_order(topic.get("commands_mentioned", [])),
                "procedures": _unique_preserve_order(topic.get("procedures", [])),
            }
        )

    return {
        "course_name": course_name or "Imported Course",
        "topics": final_topics,
        "_generation_mode": "ai",
        "_chunk_count": len(chunks),
        "_chunks_with_topics": chunks_with_topics,
        "_chunk_failures": chunk_failures,
    }


async def enrich_content(text: str) -> Dict[str, Any]:
    """Expand procedural/bulleted source material into richer educational prose."""
    client = get_client()
    user_prompt = f"""Convert this course material into rich educational content that explains every concept, command, and procedure in depth.

Original material:
{text}

Write detailed educational content covering every topic in the material."""

    try:
        response = await asyncio.wait_for(
            client.call_model("enrichment", ENRICHMENT_SYSTEM, user_prompt, max_tokens=7000),
            timeout=80,
        )
        enriched = (response or "").strip()
        if not enriched:
            raise ValueError("Enrichment returned empty content")

        return {
            "text": enriched,
            "_generation_mode": "ai",
        }
    except Exception as e:
        if 'rate-limited' in str(e).lower() or '429' in str(e).lower():
            raise
        logger.warning("Content enrichment failed, using raw text: %s", e)
        return {
            "text": text,
            "_generation_mode": "fallback",
            "_fallback_reason": str(e),
        }


async def generate_challenges(topic_data: Dict[str, Any]) -> Dict[str, Any]:
    """Generate challenges for a topic using DeepSeek V3.

    Returns a dict payload with:
    - challenges: list[dict]
    - skipped_topics: list[dict]
    """
    client = get_client()
    
    topic_name = topic_data.get("name", "Unknown Topic")
    key_concepts = topic_data.get("key_concepts", [])
    tools = topic_data.get("tools_mentioned", [])
    commands = topic_data.get("commands_mentioned", [])
    procedures = topic_data.get("procedures", [])
    domain = _detect_topic_domain(topic_data)

    topic_skip_reason = _classify_topic_sandboxability(topic_data)
    if topic_skip_reason:
        return {
            "challenges": [],
            "skipped_topics": [{"name": topic_name, "reason": topic_skip_reason}],
            "_generation_mode": "skipped",
        }
    
    system_prompt = CHALLENGE_GENERATION_SYSTEM
    challenge_type_requirements = _build_challenge_type_requirements(topic_data)
    validation_patterns = _build_validation_patterns(topic_data)
    
    user_prompt = f"""Generate practical, hands-on challenges for this topic. 
Challenges must test actual command execution, not just knowledge recall.

Topic: {topic_name}
Detected domain: {domain}
Key concepts: {key_concepts}
Tools: {tools}
Commands: {commands}
Procedures: {procedures}

{QUALITY_RULES}

{validation_patterns}

{challenge_type_requirements}

For each challenge, determine:
- type: "command" (student runs a command), "file" (student creates/modifies a file), or "output" (student runs command and output is checked)
- sandbox_image: one of ["rocky9-base", "ubuntu-wazuh", "kali-base"] — pick what makes sense
- validation_script: a bash script that exits 0 if correct, exits 1 if wrong. This runs inside the container after the student submits.

Return this exact JSON:
{{
    "challenges": [
        {{
            "question": "Clear, specific instruction",
            "hint": "Optional hint or null",
            "type": "command",
            "sandbox_type": "single",
            "sandbox_image": "rocky9-base",
            "difficulty": "easy",
            "validation_script": "#!/bin/bash\\n...",
            "expected_output": null,
            "skipped_reason": null
        }}
    ],
    "skipped_topics": [
        {{
            "name": "Topic or challenge that cannot be sandboxed",
            "reason": "multi-host or external-dependency reason"
        }}
    ]
}}

If a challenge is not sandboxable, add it to skipped_topics with a reason.
Never add non-sandboxable challenges to the challenges array.
Never replace non-sandboxable challenges with writing notes/markdown/checklists or explanation-only tasks.
Generate between 3 and 8 sandboxable challenges per topic. Make them progressively harder."""
    
    try:
        response = await asyncio.wait_for(
            client.call_model("challenge_gen", system_prompt, user_prompt, max_tokens=2400),
            timeout=20,
        )
        data = parse_json_response(response)
        if not isinstance(data, dict):
            data = {"challenges": [], "skipped_topics": []}

        challenges = data.get("challenges", []) if isinstance(data.get("challenges", []), list) else []
        skipped_topics = data.get("skipped_topics", []) if isinstance(data.get("skipped_topics", []), list) else []

        # Enforce single-container constraints defensively even if model output drifts.
        filtered_challenges: List[Dict[str, Any]] = []
        seen_local_fingerprints: set[str] = set()
        for challenge in challenges:
            if not isinstance(challenge, dict):
                continue

            question = str(challenge.get("question") or "")
            validation_script = str(challenge.get("validation_script") or "")
            sandbox_type = str(challenge.get("sandbox_type") or "single").strip().lower()
            if sandbox_type != "single":
                skipped_topics.append(
                    {
                        "name": question or "Unnamed challenge",
                        "reason": challenge.get("skipped_reason") or f"{sandbox_type}: not sandboxable in single-container Docker",
                    }
                )
                continue

            filler_reason = _is_generic_filler_challenge(question, validation_script)
            if filler_reason:
                skipped_topics.append(
                    {
                        "name": question or "Unnamed challenge",
                        "reason": f"quality-filter: {filler_reason}",
                    }
                )
                continue

            domain_quality_reason = _is_low_value_domain_challenge(question, validation_script, topic_data)
            if domain_quality_reason:
                skipped_topics.append(
                    {
                        "name": question or "Unnamed challenge",
                        "reason": f"quality-filter: {domain_quality_reason}",
                    }
                )
                continue

            weak_script_reason = _is_weak_validation_script(validation_script)
            if weak_script_reason:
                skipped_topics.append(
                    {
                        "name": question or "Unnamed challenge",
                        "reason": f"quality-filter: {weak_script_reason}",
                    }
                )
                continue

            local_fingerprint = _challenge_fingerprint(question, validation_script)
            if local_fingerprint in seen_local_fingerprints:
                skipped_topics.append(
                    {
                        "name": question or "Unnamed challenge",
                        "reason": "quality-filter: duplicate template in same topic",
                    }
                )
                continue

            seen_local_fingerprints.add(local_fingerprint)

            challenge.setdefault("sandbox_type", "single")
            challenge.setdefault("skipped_reason", None)
            filtered_challenges.append(challenge)

        if not filtered_challenges:
            logger.warning("AI challenge generation returned empty list; using fallback challenges")
            return {
                "challenges": _fallback_generate_challenges(topic_data),
                "skipped_topics": skipped_topics,
                "_generation_mode": "fallback",
                "_fallback_reason": "AI returned no sandboxable challenges",
            }

        domain_mix_reason = _validate_domain_challenge_mix(filtered_challenges, topic_data)
        if domain_mix_reason:
            skipped_topics.append(
                {
                    "name": topic_name,
                    "reason": f"quality-filter: {domain_mix_reason}",
                }
            )
            fallback_challenges = _fallback_generate_challenges(topic_data)
            if fallback_challenges:
                return {
                    "challenges": fallback_challenges,
                    "skipped_topics": skipped_topics,
                    "_generation_mode": "fallback",
                    "_fallback_reason": domain_mix_reason,
                }
            return {
                "challenges": [],
                "skipped_topics": skipped_topics,
                "_generation_mode": "skipped",
                "_fallback_reason": domain_mix_reason,
            }

        return {
            "challenges": filtered_challenges,
            "skipped_topics": skipped_topics,
            "_generation_mode": "ai",
        }
    except Exception as e:
        if 'rate-limited' in str(e).lower() or '429' in str(e).lower():
            raise
        logger.warning(f"Failed to generate challenges via AI, using fallback challenges: {e}")
        return {
            "challenges": _fallback_generate_challenges(topic_data),
            "skipped_topics": [],
            "_generation_mode": "fallback",
            "_fallback_reason": str(e),
        }


async def review_validation_script(question: str, script: str, sandbox_image: str) -> Dict[str, Any]:
    """Review validation script using Qwen3 Coder."""
    if not ENABLE_AI_REVIEW_PASSES:
        return {"valid": True, "issues": None, "fixed_script": None}

    client = get_client()
    
    system_prompt = """You are a bash script reviewer. Return ONLY valid JSON."""
    
    user_prompt = f"""Review this bash validation script for a lab challenge.
The script runs inside a Docker container after a student completes a task.
It should exit 0 on success, exit 1 on failure.

Question: {question}
Script: {script}
Container image: {sandbox_image}

Is this script correct and will it reliably validate the student's work?
Return JSON: {{"valid": true/false, "issues": "description or null", "fixed_script": "corrected script or null"}}"""
    
    try:
        response = await asyncio.wait_for(
            client.call_model("validator_review", system_prompt, user_prompt, max_tokens=2000),
            timeout=12,
        )
        return parse_json_response(response)
    except Exception as e:
        logger.error(f"Failed to review validation script: {e}")
        return {"valid": True, "issues": None, "fixed_script": None}


def _local_sanity_review_challenge(question: str, validation_script: str, topic_name: Optional[str] = None) -> Dict[str, Any]:
    """Fast local sanity checks to block obviously broken challenges without AI dependency."""
    q = (question or "").strip()
    q_lower = q.lower()
    script = (validation_script or "").strip()

    if not q:
        return {"approved": False, "reason": "Empty challenge question", "improved_question": None}

    filler_reason = _is_generic_filler_challenge(q, script)
    if filler_reason:
        return {
            "approved": False,
            "reason": f"Low-value generic template ({filler_reason})",
            "improved_question": None,
        }

    topic_hint = {"name": topic_name or ""}
    domain_low_value_reason = _is_low_value_domain_challenge(q, script, topic_hint)
    if domain_low_value_reason:
        return {
            "approved": False,
            "reason": f"Low-value domain template ({domain_low_value_reason})",
            "improved_question": None,
        }

    if "verify your environment" in q_lower:
        return {
            "approved": False,
            "reason": "Question is too vague ('verify your environment')",
            "improved_question": None,
        }

    hypervisor_terms = ["vmware", "virtualbox", "hyper-v", "vcenter", "esxi", "vmnet", "ova", "ovf"]
    if any(term in q_lower for term in hypervisor_terms):
        return {
            "approved": False,
            "reason": "Requires hypervisor/VM tooling outside single-container Docker sandbox",
            "improved_question": None,
        }

    if "snapshot" in q_lower and (" vm" in q_lower or "virtual" in q_lower or "ova" in q_lower or "ovf" in q_lower):
        return {
            "approved": False,
            "reason": "Snapshot operation appears to be VM/hypervisor action, not a container CLI command",
            "improved_question": None,
        }

    if any(term in q_lower for term in ["github", "gitlab", "bitbucket", "push to repo", "pull request"]):
        return {
            "approved": False,
            "reason": "Requires external dependency/service not guaranteed in sandbox",
            "improved_question": None,
        }

    if any(term in q_lower for term in ["ssh to", "remote host", "managed node", "second host", "another machine"]):
        return {
            "approved": False,
            "reason": "Requires multi-host networking beyond single-container sandbox",
            "improved_question": None,
        }

    weak_script_reason = _is_weak_validation_script(script)
    if weak_script_reason:
        return {
            "approved": False,
            "reason": f"Weak validation script ({weak_script_reason})",
            "improved_question": None,
        }

    return {"approved": True, "reason": None, "improved_question": None}


async def sanity_review_challenge(
    question: str,
    validation_script: str,
    sandbox_image: str,
    topic_name: Optional[str] = None,
) -> Dict[str, Any]:
    """AI sanity review pass to reject non-sandboxable or low-quality challenges."""
    local = _local_sanity_review_challenge(question, validation_script, topic_name)
    if not local.get("approved"):
        return local

    if not ENABLE_AI_REVIEW_PASSES:
        return local

    # Skip external AI in test environment to keep tests deterministic/offline.
    settings = get_settings()
    api_key = (settings.OPENROUTER_API_KEY or "").strip()
    if api_key.startswith("test-"):
        return local

    client = get_client()
    user_prompt = f"""Review this generated challenge for Docker-sandbox viability.

Question: {question}
Validation script:
{validation_script}
Sandbox image: {sandbox_image}

Return JSON exactly in this format:
{{
  "approved": true/false,
  "reason": "why rejected or null if approved",
  "improved_question": "better version of the question or null"
}}"""

    try:
        response = await asyncio.wait_for(
            client.call_model("sanity_review", CHALLENGE_SANITY_REVIEW_SYSTEM, user_prompt, max_tokens=1200),
            timeout=15,
        )
        parsed = parse_json_response(response)
        approved = bool(parsed.get("approved"))
        reason = parsed.get("reason")
        improved_question = parsed.get("improved_question")

        return {
            "approved": approved,
            "reason": reason if isinstance(reason, str) or reason is None else str(reason),
            "improved_question": improved_question if isinstance(improved_question, str) or improved_question is None else None,
        }
    except Exception as e:
        logger.warning(f"Sanity review unavailable; keeping challenge via local checks: {e}")
        return local


def save_challenge_cache(course_id: str, topic_slug: str, challenges: List[Dict[str, Any]]) -> None:
    """Save challenge cache to disk."""
    CHALLENGES_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    course_dir = CHALLENGES_CACHE_DIR / course_id
    course_dir.mkdir(exist_ok=True)

    raw = (topic_slug or "unknown-topic").strip().lower()
    normalized = raw.replace("_", "-")
    normalized = re.sub(r"[\\/]+", "-", normalized)
    normalized = re.sub(r"[^a-z0-9-]+", "-", normalized)
    normalized = re.sub(r"-{2,}", "-", normalized).strip("-")
    if not normalized:
        normalized = "topic"

    # Add short hash so sanitized/truncated names stay unique and deterministic.
    topic_key = f"{normalized[:40]}-{hashlib.sha1(raw.encode('utf-8')).hexdigest()[:8]}"

    cache_file = course_dir / f"{topic_key}.json"
    cache_file.parent.mkdir(parents=True, exist_ok=True)
    with open(cache_file, "w") as f:
        json.dump(challenges, f, indent=2)


def load_challenge_cache(course_id: str, topic_slug: str) -> Optional[List[Dict[str, Any]]]:
    """Load challenge cache from disk."""
    raw = (topic_slug or "unknown-topic").strip().lower()
    normalized = raw.replace("_", "-")
    normalized = re.sub(r"[\\/]+", "-", normalized)
    normalized = re.sub(r"[^a-z0-9-]+", "-", normalized)
    normalized = re.sub(r"-{2,}", "-", normalized).strip("-")
    if not normalized:
        normalized = "topic"

    topic_key = f"{normalized[:40]}-{hashlib.sha1(raw.encode('utf-8')).hexdigest()[:8]}"
    cache_file = CHALLENGES_CACHE_DIR / course_id / f"{topic_key}.json"
    if cache_file.exists():
        with open(cache_file) as f:
            return json.load(f)
    return None


def _update_job_progress(job_id: str, percent: int) -> None:
    """Update job progress in database."""
    with Session(engine) as session:
        job = session.get(ImportJob, job_id)
        if job:
            job.update_progress(percent)
            session.commit()


def _add_job_log(job_id: str, message: str, level: str = "info") -> None:
    """Add log entry to job."""
    with Session(engine) as session:
        job = session.get(ImportJob, job_id)
        if job:
            job.add_log(message, level)
            session.commit()


def _update_job_status(job_id: str, status: str, course_id: Optional[str] = None, 
                        error_message: Optional[str] = None) -> None:
    """Update job status in database."""
    with Session(engine) as session:
        job = session.get(ImportJob, job_id)
        if job:
            job.status = status
            if course_id:
                job.course_id = course_id
            if error_message:
                job.error_message = error_message
            if status in ["completed", "error"]:
                job.completed_at = datetime.utcnow()
            job.updated_at = datetime.utcnow()
            session.commit()


def _update_job_counts(job_id: str, topics_count: int, challenges_count: int) -> None:
    """Update job counts in database."""
    with Session(engine) as session:
        job = session.get(ImportJob, job_id)
        if job:
            job.topics_count = topics_count
            job.challenges_count = challenges_count
            session.commit()


def _is_job_cancelled(job_id: str) -> bool:
    """Check whether a job was cancelled by the user."""
    with Session(engine) as session:
        job = session.get(ImportJob, job_id)
        return bool(job and job.status == "error" and job.error_message == "Cancelled by user")


async def _generate_topic_challenge_batches(topics: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Generate challenges for topics with bounded concurrency.

    Returns a list aligned with input order containing:
    - generated: optional dict payload from generate_challenges
    - error: optional Exception
    """
    semaphore = asyncio.Semaphore(CHALLENGE_GEN_BATCH_CONCURRENCY)

    async def _run_one(idx: int, topic_data: Dict[str, Any]):
        topic_name = str(topic_data.get("name") or f"Topic {idx + 1}")
        try:
            async with semaphore:
                generated = await generate_challenges(topic_data)
            return idx, {"generated": generated, "error": None}
        except Exception as e:
            return idx, {"generated": None, "error": e, "topic_name": topic_name}

    tasks = [_run_one(idx, topic_data) for idx, topic_data in enumerate(topics)]
    gathered = await asyncio.gather(*tasks)

    ordered: List[Dict[str, Any]] = [{"generated": None, "error": None} for _ in topics]
    for idx, payload in gathered:
        ordered[idx] = payload
    return ordered


async def process_file_with_job(job_id: str, file_path: str, source_filename: Optional[str] = None) -> None:
    """Process a file with job tracking and progress updates."""
    async with _job_semaphore:
        try:
            display_source_file = source_filename or Path(file_path).name
            _update_job_status(job_id, "processing")
            _add_job_log(job_id, f"Starting processing of {display_source_file}")
            
            # Step 1: Compute source hash and check for duplicates (0-5%)
            _add_job_log(job_id, "Computing file hash for duplicate detection...")
            source_hash = compute_source_hash(file_path)
            
            with Session(engine) as session:
                existing = session.exec(select(Course).where(Course.source_hash == source_hash)).first()
                if existing:
                    # Check if course has content - if not, delete and reprocess
                    if existing.topic_count == 0 or existing.challenge_count == 0:
                        _add_job_log(job_id, f"Found empty course '{existing.title}' - deleting and reprocessing", "warn")
                        session.delete(existing)
                        session.commit()
                        _add_job_log(job_id, "Deleted empty course, starting fresh processing...")
                    else:
                        _add_job_log(job_id, f"File already processed as course: {existing.title}", "info")
                        _update_job_counts(job_id, existing.topic_count, existing.challenge_count)
                        _update_job_status(job_id, "completed", course_id=existing.id)
                        _update_job_progress(job_id, 100)
                        return
            
            _update_job_progress(job_id, 5)
            
            # Step 2: Parse file (5-15%)
            if _is_job_cancelled(job_id):
                _add_job_log(job_id, "Job cancelled before parse step", "warn")
                return

            _add_job_log(job_id, "Parsing document...")
            if file_path.lower().endswith(".pdf"):
                text = parse_pdf(file_path)
            elif file_path.lower().endswith(".pptx"):
                text = parse_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_path}")

            if not text or not text.strip():
                raise ValueError("No readable text could be extracted from the uploaded document")
            
            _add_job_log(job_id, f"Extracted {len(text)} characters from document")
            raw_chunks = chunk_text(text, max_tokens=MAX_CHUNK_TOKENS)
            _add_job_log(job_id, f"Raw content chunking: {len(raw_chunks)} chunks")
            _update_job_progress(job_id, 15)

            # Step 2.5: Runtime enrichment disabled (keeps extraction grounded + faster)
            enriched_text = text
            if ENABLE_RUNTIME_ENRICHMENT:
                if _is_job_cancelled(job_id):
                    _add_job_log(job_id, "Job cancelled before enrichment", "warn")
                    return

                _add_job_log(job_id, "Enriching content with AI...")
                enriched_payload = await enrich_content(text)
                enriched_text = str(enriched_payload.get("text") or text)
                if enriched_payload.get("_generation_mode") == "fallback":
                    _add_job_log(
                        job_id,
                        f"Content enrichment unavailable; using raw text ({enriched_payload.get('_fallback_reason', 'no reason')})",
                        "warn",
                    )
                _add_job_log(job_id, f"Content enriched: {len(text)} → {len(enriched_text)} chars")
                enriched_chunks = chunk_text(enriched_text, max_tokens=MAX_CHUNK_TOKENS)
                _add_job_log(job_id, f"Enriched content chunking: {len(enriched_chunks)} chunks")
            else:
                _add_job_log(job_id, "Enrichment stage disabled; using raw extracted text")
            
            # Step 3: Extract topics (15-40%)
            if _is_job_cancelled(job_id):
                _add_job_log(job_id, "Job cancelled before topic extraction", "warn")
                return

            _add_job_log(job_id, "Extracting topics with AI...")
            topics_data = await extract_topics(enriched_text)
            if topics_data.get("_generation_mode") == "fallback":
                _add_job_log(
                    job_id,
                    f"AI topic extraction unavailable; used fallback extraction ({topics_data.get('_fallback_reason', 'no reason')})",
                    "warn",
                )
            if topics_data.get("_chunk_count"):
                _add_job_log(
                    job_id,
                    (
                        f"Topic extraction chunk stats: "
                        f"{topics_data.get('_chunks_with_topics', 0)}/{topics_data.get('_chunk_count', 0)} with topics, "
                        f"failures={topics_data.get('_chunk_failures', 0)}"
                    ),
                )
            course_name = topics_data.get("course_name", "Unknown Course")
            topics = topics_data.get("topics", [])

            if not topics:
                raise ValueError("No topics were extracted from the document")
            
            _add_job_log(job_id, f"Found {len(topics)} topics in course: {course_name}")
            _update_job_progress(job_id, 40)
            
            # Step 4: Create course
            course_id = str(uuid.uuid4())
            with Session(engine) as session:
                course = Course(
                    id=course_id,
                    title=course_name,
                    description=f"Auto-generated from {display_source_file}",
                    source_file=display_source_file,
                    source_hash=source_hash,
                )
                session.add(course)
                session.commit()
            
            _add_job_log(job_id, f"Created course: {course_name}")
            
            # Step 5: Process each topic (40-90%)
            total_challenges = 0
            retained_topics_count = 0
            seen_challenge_fingerprints: set[str] = set()
            topic_progress_increment = 50 / len(topics) if topics else 50

            _add_job_log(
                job_id,
                f"Generating challenge drafts for {len(topics)} topics (concurrency={CHALLENGE_GEN_BATCH_CONCURRENCY})...",
            )
            batched_topic_generation = await _generate_topic_challenge_batches(topics)
            _update_job_progress(job_id, 45)
            
            for idx, topic_data in enumerate(topics):
                if _is_job_cancelled(job_id):
                    _add_job_log(job_id, "Job cancelled during topic processing", "warn")
                    return

                topic_name = topic_data.get("name", "Unknown Topic")
                topic_order = topic_data.get("order", idx + 1)
                
                _add_job_log(job_id, f"Processing topic {idx + 1}/{len(topics)}: {topic_name}")
                topic_start_progress = 40 + int(idx * topic_progress_increment)
                topic_end_progress = 40 + int((idx + 1) * topic_progress_increment)
                _update_job_progress(job_id, min(topic_start_progress + 5, 89))
                
                # Create topic
                topic_id = str(uuid.uuid4())
                with Session(engine) as session:
                    topic = Topic(
                        id=topic_id,
                        course_id=course_id,
                        name=topic_name,
                        order=topic_order,
                    )
                    session.add(topic)
                    session.commit()
                
                # Use pre-generated batch result for this topic.
                _add_job_log(job_id, f"Preparing challenges for topic: {topic_name}")
                batch_payload = batched_topic_generation[idx] if idx < len(batched_topic_generation) else {"generated": None, "error": None}
                batch_error = batch_payload.get("error")
                if batch_error:
                    _add_job_log(job_id, f"Failed to generate challenges for {topic_name}: {batch_error}", "error")
                    continue

                generated = batch_payload.get("generated")
                if isinstance(generated, dict):
                    challenges_data = generated.get("challenges", []) or []
                    skipped = generated.get("skipped_topics", []) or []
                    generation_mode = str(generated.get("_generation_mode") or "ai")
                    if generation_mode == "fallback":
                        _add_job_log(
                            job_id,
                            f"AI challenge generation unavailable; used fallback challenges ({generated.get('_fallback_reason', 'no reason')})",
                            "warn",
                        )
                    elif generation_mode == "skipped":
                        _add_job_log(job_id, f"Topic pre-classified non-sandboxable: {topic_name}", "warn")
                else:
                    # Backwards compatibility for legacy monkeypatch/tests returning list.
                    challenges_data = generated or []
                    skipped = []

                if skipped:
                    for skip in skipped:
                        if not isinstance(skip, dict):
                            continue
                        logger.info(
                            "[grinder] Skipped non-sandboxable challenge: %s — %s",
                            skip.get("name", "Unnamed"),
                            skip.get("reason", "unspecified"),
                        )
                    _add_job_log(job_id, f"Skipped {len(skipped)} multi-host/external challenges (not sandboxable in Docker)")

                _add_job_log(job_id, f"Generated {len(challenges_data)} challenges for {topic_name}")
                _update_job_progress(job_id, min(topic_start_progress + 20, max(topic_start_progress + 5, topic_end_progress - 1)))
                
                # Process each challenge
                total_topic_challenges = len(challenges_data)
                approved_topic_challenges = 0
                topic_fingerprints: set[str] = set()
                for challenge_idx, challenge_data in enumerate(challenges_data):
                    if _is_job_cancelled(job_id):
                        _add_job_log(job_id, "Job cancelled during challenge processing", "warn")
                        return

                    question = challenge_data.get("question", "")
                    validation_script = challenge_data.get("validation_script", "")
                    
                    # Review validation script
                    try:
                        review = await review_validation_script(
                            question, validation_script, challenge_data.get("sandbox_image", "rocky9-base")
                        )
                        if review.get("valid") and review.get("fixed_script"):
                            validation_script = review["fixed_script"]
                    except Exception as e:
                        _add_job_log(job_id, f"Validation script review failed: {e}", "warn")

                    # Sanity review to reject non-sandboxable or low-quality challenges.
                    try:
                        sanity = await sanity_review_challenge(
                            question,
                            validation_script,
                            challenge_data.get("sandbox_image", "rocky9-base"),
                            topic_name,
                        )
                    except Exception as e:
                        _add_job_log(job_id, f"Sanity review failed (defaulting to keep): {e}", "warn")
                        sanity = {"approved": True, "reason": None, "improved_question": None}

                    if not sanity.get("approved"):
                        reject_reason = sanity.get("reason") or "rejected by sanity review"
                        logger.info("[grinder] Challenge rejected by sanity review: %s", reject_reason)
                        _add_job_log(job_id, f"Rejected: '{question[:80]}' — {reject_reason}", "warn")
                        continue

                    improved_question = sanity.get("improved_question")
                    if isinstance(improved_question, str) and improved_question.strip():
                        question = improved_question.strip()

                    fingerprint = _challenge_fingerprint(question, validation_script)
                    if fingerprint in seen_challenge_fingerprints:
                        _add_job_log(job_id, f"Skipped duplicate template across topics: '{question[:80]}'", "warn")
                        continue

                    # Compute deterministic challenge ID after all challenge modifications.
                    challenge_id = compute_challenge_id(question, validation_script, topic_id)

                    # Check if challenge already exists
                    with Session(engine) as session:
                        existing_challenge = session.get(Challenge, challenge_id)
                        if existing_challenge:
                            continue
                    
                    # Create challenge
                    with Session(engine) as session:
                        challenge = Challenge(
                            id=challenge_id,
                            course_id=course_id,
                            topic_id=topic_id,
                            type=challenge_data.get("type", "command"),
                            question=question,
                            hint=challenge_data.get("hint"),
                            sandbox_image=challenge_data.get("sandbox_image", "rocky9-base"),
                            validation_script=validation_script,
                            expected_output=challenge_data.get("expected_output"),
                            difficulty=challenge_data.get("difficulty", "easy"),
                            order=challenge_data.get("order", 0),
                        )
                        session.add(challenge)
                        session.commit()
                        total_challenges += 1
                        approved_topic_challenges += 1
                        seen_challenge_fingerprints.add(fingerprint)
                        topic_fingerprints.add(fingerprint)

                    # Granular per-challenge progress (helps avoid "stuck at 40%" UI feel)
                    if total_topic_challenges > 0:
                        challenge_band_start = topic_start_progress + 20
                        challenge_band_span = max(1, (topic_end_progress - challenge_band_start - 1))
                        challenge_progress = challenge_band_start + int(((challenge_idx + 1) / total_topic_challenges) * challenge_band_span)
                        _update_job_progress(job_id, min(challenge_progress, topic_end_progress - 1))
                
                if approved_topic_challenges < MIN_APPROVED_CHALLENGES_PER_TOPIC:
                    with Session(engine) as session:
                        session.exec(delete(Challenge).where(Challenge.topic_id == topic_id))
                        session.exec(delete(Topic).where(Topic.id == topic_id))
                        session.commit()

                    total_challenges -= approved_topic_challenges
                    seen_challenge_fingerprints.difference_update(topic_fingerprints)
                    _add_job_log(
                        job_id,
                        (
                            f"Skipped topic '{topic_name}' after quality gate: "
                            f"{approved_topic_challenges} approved challenges (minimum {MIN_APPROVED_CHALLENGES_PER_TOPIC})"
                        ),
                        "warn",
                    )
                else:
                    retained_topics_count += 1
                    # Save cache only for retained topics.
                    save_challenge_cache(course_id, topic_name, challenges_data)
                
                # Update progress per topic
                current_progress = topic_end_progress
                _update_job_progress(job_id, min(current_progress, 90))
            
            # Step 6: Update course counts (90-100%)
            if retained_topics_count == 0 or total_challenges == 0:
                with Session(engine) as session:
                    session.exec(delete(Challenge).where(Challenge.course_id == course_id))
                    session.exec(delete(Topic).where(Topic.course_id == course_id))
                    session.exec(delete(Course).where(Course.id == course_id))
                    session.commit()

                skipped_topics_count = max(0, len(topics) - retained_topics_count)
                raise ValueError(
                    (
                        "No publishable topics remained after quality filtering "
                        f"({retained_topics_count} retained, {skipped_topics_count} skipped, {total_challenges} challenges). "
                        "Likely caused by temporary AI provider limits or low-extractability source content. "
                        "Please retry later or review grinder logs."
                    )
                )

            with Session(engine) as session:
                course = session.get(Course, course_id)
                if course:
                    course.topic_count = retained_topics_count
                    course.challenge_count = total_challenges
                    session.commit()
            
            _update_job_counts(job_id, retained_topics_count, total_challenges)
            skipped_topics_count = max(0, len(topics) - retained_topics_count)
            _add_job_log(
                job_id,
                f"Processing complete: {retained_topics_count} retained topics, {skipped_topics_count} skipped topics, {total_challenges} challenges",
            )
            _update_job_progress(job_id, 100)
            if not _is_job_cancelled(job_id):
                _update_job_status(job_id, "completed", course_id=course_id)
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"Error processing job {job_id}: {error_msg}")
            _add_job_log(job_id, f"Processing failed: {error_msg}", "error")
            _update_job_status(job_id, "error", error_message=error_msg)
            raise


async def process_file(file_path: str, source_filename: Optional[str] = None) -> Dict[str, Any]:
    """Main entry point: process a course file end-to-end (sync version for backwards compatibility)."""
    display_source_file = source_filename or Path(file_path).name
    logger.info(f"Processing file: {display_source_file} (path: {file_path})")
    
    # Step 1: Compute source hash and check for duplicates
    source_hash = compute_source_hash(file_path)
    
    with Session(engine) as session:
        existing = session.exec(select(Course).where(Course.source_hash == source_hash)).first()
        if existing:
            # Check if course has content - if not, delete and reprocess
            if existing.topic_count == 0 or existing.challenge_count == 0:
                logger.info(f"Found empty course '{existing.title}' - deleting and reprocessing")
                session.delete(existing)
                session.commit()
                logger.info("Deleted empty course, starting fresh processing...")
            else:
                logger.info(f"File already processed: {existing.id}")
                return {
                    "course_id": existing.id,
                    "topics_count": existing.topic_count,
                    "challenges_count": existing.challenge_count,
                    "status": "already_processed",
                }
    
    # Step 2: Parse file
    if file_path.lower().endswith(".pdf"):
        text = parse_pdf(file_path)
    elif file_path.lower().endswith(".pptx"):
        text = parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

    if not text or not text.strip():
        raise ValueError("No readable text could be extracted from the uploaded document")

    logger.info("Raw content length: %s chars, chunks: %s", len(text), len(chunk_text(text, max_tokens=MAX_CHUNK_TOKENS)))

    # Step 2.5: Runtime enrichment disabled by default.
    enriched_text = text
    if ENABLE_RUNTIME_ENRICHMENT:
        logger.info("Enriching content with AI...")
        enriched_payload = await enrich_content(text)
        enriched_text = str(enriched_payload.get("text") or text)
        if enriched_payload.get("_generation_mode") == "fallback":
            logger.warning(
                "[grinder] Content enrichment unavailable; using raw text (%s)",
                enriched_payload.get("_fallback_reason", "no reason"),
            )
        logger.info("Content enriched: %s -> %s chars", len(text), len(enriched_text))
        logger.info("Enriched content chunks: %s", len(chunk_text(enriched_text, max_tokens=MAX_CHUNK_TOKENS)))
    else:
        logger.info("Enrichment stage disabled; using raw extracted text")
    
    # Step 3: Extract topics
    logger.info("Extracting topics...")
    topics_data = await extract_topics(enriched_text)
    if topics_data.get("_generation_mode") == "fallback":
        logger.warning(
            "[grinder] AI topic extraction unavailable; used fallback extraction (%s)",
            topics_data.get("_fallback_reason", "no reason"),
        )
    if topics_data.get("_chunk_count"):
        logger.info(
            "[grinder] Topic extraction chunk stats: %s/%s with topics, failures=%s",
            topics_data.get("_chunks_with_topics", 0),
            topics_data.get("_chunk_count", 0),
            topics_data.get("_chunk_failures", 0),
        )
    course_name = topics_data.get("course_name", "Unknown Course")
    topics = topics_data.get("topics", [])

    if not topics:
        raise ValueError("No topics were extracted from the document")
    
    # Step 4: Create course
    course_id = str(uuid.uuid4())
    with Session(engine) as session:
        course = Course(
            id=course_id,
            title=course_name,
            description=f"Auto-generated from {display_source_file}",
            source_file=display_source_file,
            source_hash=source_hash,
        )
        session.add(course)
        session.commit()
    
    # Step 5: Process each topic
    total_challenges = 0
    retained_topics_count = 0
    seen_challenge_fingerprints: set[str] = set()
    logger.info(
        "Generating challenge drafts for %s topics (concurrency=%s)",
        len(topics),
        CHALLENGE_GEN_BATCH_CONCURRENCY,
    )
    batched_topic_generation = await _generate_topic_challenge_batches(topics)
    for topic_index, topic_data in enumerate(topics):
        topic_name = topic_data.get("name", "Unknown Topic")
        topic_order = topic_data.get("order", 0)
        
        logger.info(f"Processing topic: {topic_name}")
        
        # Create topic
        topic_id = str(uuid.uuid4())
        with Session(engine) as session:
            topic = Topic(
                id=topic_id,
                course_id=course_id,
                name=topic_name,
                order=topic_order,
            )
            session.add(topic)
            session.commit()
        
        # Use pre-generated batch result for this topic.
        batch_payload = batched_topic_generation[topic_index] if topic_index < len(batched_topic_generation) else {"generated": None, "error": None}
        batch_error = batch_payload.get("error")
        if batch_error:
            logger.error("Failed to generate challenges for %s: %s", topic_name, batch_error)
            continue

        generated = batch_payload.get("generated")
        if isinstance(generated, dict):
            challenges_data = generated.get("challenges", []) or []
            skipped = generated.get("skipped_topics", []) or []
            if generated.get("_generation_mode") == "fallback":
                logger.warning(
                    "[grinder] AI challenge generation unavailable; used fallback challenges (%s)",
                    generated.get("_fallback_reason", "no reason"),
                )
        else:
            # Backwards compatibility for legacy monkeypatch/tests returning list.
            challenges_data = generated or []
            skipped = []

        if skipped:
            for skip in skipped:
                if not isinstance(skip, dict):
                    continue
                logger.info(
                    "[grinder] Skipped non-sandboxable challenge: %s — %s",
                    skip.get("name", "Unnamed"),
                    skip.get("reason", "unspecified"),
                )
            logger.info("[grinder] Skipped %d multi-host/external challenges (not sandboxable in Docker)", len(skipped))
        
        # Process each challenge
        approved_topic_challenges = 0
        topic_fingerprints: set[str] = set()
        for challenge_data in challenges_data:
            question = challenge_data.get("question", "")
            validation_script = challenge_data.get("validation_script", "")
            
            # Review validation script
            try:
                review = await review_validation_script(
                    question, validation_script, challenge_data.get("sandbox_image", "rocky9-base")
                )
                if review.get("valid") and review.get("fixed_script"):
                    validation_script = review["fixed_script"]
            except Exception as e:
                logger.warning(f"Validation script review failed: {e}")

            # Sanity review to reject non-sandboxable or low-quality challenges.
            try:
                sanity = await sanity_review_challenge(
                    question,
                    validation_script,
                    challenge_data.get("sandbox_image", "rocky9-base"),
                    topic_name,
                )
            except Exception as e:
                logger.warning(f"Sanity review failed (defaulting to keep): {e}")
                sanity = {"approved": True, "reason": None, "improved_question": None}

            if not sanity.get("approved"):
                reject_reason = sanity.get("reason") or "rejected by sanity review"
                logger.info("[grinder] Challenge rejected by sanity review: %s", reject_reason)
                continue

            improved_question = sanity.get("improved_question")
            if isinstance(improved_question, str) and improved_question.strip():
                question = improved_question.strip()

            fingerprint = _challenge_fingerprint(question, validation_script)
            if fingerprint in seen_challenge_fingerprints:
                logger.info("[grinder] Skipping duplicate template across topics: %s", question[:80])
                continue

            # Compute deterministic challenge ID after all challenge modifications.
            challenge_id = compute_challenge_id(question, validation_script, topic_id)

            # Check if challenge already exists
            with Session(engine) as session:
                existing_challenge = session.get(Challenge, challenge_id)
                if existing_challenge:
                    logger.info(f"Challenge already exists: {challenge_id[:16]}...")
                    continue
            
            # Create challenge
            with Session(engine) as session:
                challenge = Challenge(
                    id=challenge_id,
                    course_id=course_id,
                    topic_id=topic_id,
                    type=challenge_data.get("type", "command"),
                    question=question,
                    hint=challenge_data.get("hint"),
                    sandbox_image=challenge_data.get("sandbox_image", "rocky9-base"),
                    validation_script=validation_script,
                    expected_output=challenge_data.get("expected_output"),
                    difficulty=challenge_data.get("difficulty", "easy"),
                    order=challenge_data.get("order", 0),
                )
                session.add(challenge)
                session.commit()
                total_challenges += 1
                approved_topic_challenges += 1
                seen_challenge_fingerprints.add(fingerprint)
                topic_fingerprints.add(fingerprint)

        if approved_topic_challenges < MIN_APPROVED_CHALLENGES_PER_TOPIC:
            with Session(engine) as session:
                session.exec(delete(Challenge).where(Challenge.topic_id == topic_id))
                session.exec(delete(Topic).where(Topic.id == topic_id))
                session.commit()

            total_challenges -= approved_topic_challenges
            seen_challenge_fingerprints.difference_update(topic_fingerprints)
            logger.warning(
                "[grinder] Skipping topic '%s' after quality gate: %s approved (minimum %s)",
                topic_name,
                approved_topic_challenges,
                MIN_APPROVED_CHALLENGES_PER_TOPIC,
            )
            continue

        retained_topics_count += 1
        
        # Save cache
        save_challenge_cache(course_id, topic_name, challenges_data)
    
    # Update course counts
    if retained_topics_count == 0 or total_challenges == 0:
        with Session(engine) as session:
            session.exec(delete(Challenge).where(Challenge.course_id == course_id))
            session.exec(delete(Topic).where(Topic.course_id == course_id))
            session.exec(delete(Course).where(Course.id == course_id))
            session.commit()

        skipped_topics_count = max(0, len(topics) - retained_topics_count)
        raise ValueError(
            (
                "No publishable topics remained after quality filtering "
                f"({retained_topics_count} retained, {skipped_topics_count} skipped, {total_challenges} challenges). "
                "Likely caused by temporary AI provider limits or low-extractability source content."
            )
        )

    with Session(engine) as session:
        course = session.get(Course, course_id)
        if course:
            course.topic_count = retained_topics_count
            course.challenge_count = total_challenges
            session.commit()
    
    logger.info(f"Processing complete: {len(topics)} topics, {total_challenges} challenges")
    
    return {
        "course_id": course_id,
        "course_name": course_name,
        "topics_count": retained_topics_count,
        "challenges_count": total_challenges,
        "status": "success",
    }


def get_queue_status() -> Dict[str, Any]:
    """Get grinder queue status."""
    return {
        "status": "idle",
        "queue_length": 0,
        "last_processed": None,
    }
